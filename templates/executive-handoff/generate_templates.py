#!/usr/bin/env python3
"""Generate 5 Executive Handoff Summary PowerPoint templates in distinct art styles.

Shell (11 slides, all single slides):
  1. Title            - Executive Handoff Summary / [Plan Name] / Overview of the
                        Corporate Retirement Plan Process
  2. Executive Overview
  3. Project Status
  4. Project Scope
  5. Milestones
  6. Work Stream Status
  7. Major Decisions
  8. Risk Overview
  9. Stakeholder Review
 10. Immediate Priorities
 11. Lessons Learned
"""

import math
import os
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = os.path.dirname(os.path.abspath(__file__))
ART = os.path.join(HERE, "art")
OUT = os.environ.get("PPT_OUT", HERE)
os.makedirs(ART, exist_ok=True)
os.makedirs(OUT, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------- content ---

SECTIONS = [
    # (key, title, kicker, bullets)
    ("executive-overview", "Executive Overview", "Section 01",
     ["[One-paragraph executive summary of the plan transition]",
      "[Business context: why this handoff is happening now]",
      "[Headline outcomes achieved to date]",
      "[What the incoming owner needs to know first]"]),
    ("project-status", "Project Status", "Section 02",
     ["[Current status: On Track / At Risk / Off Track]",
      "[Percent complete and current phase]",
      "[Work completed this period]",
      "[Work planned for next period]"]),
    ("project-scope", "Project Scope", "Section 03",
     ["[In scope: services, plans, and populations covered]",
      "[Out of scope: explicitly excluded items]",
      "[Key assumptions and dependencies]",
      "[Approved scope changes to date]"]),
    ("milestones", "Milestones", "Section 04",
     ["[Milestone 1 — target date — status]",
      "[Milestone 2 — target date — status]",
      "[Milestone 3 — target date — status]",
      "[Milestone 4 — target date — status]"]),
    ("work-stream-status", "Work Stream Status", "Section 05",
     ["[Work stream A — owner — status — next step]",
      "[Work stream B — owner — status — next step]",
      "[Work stream C — owner — status — next step]",
      "[Work stream D — owner — status — next step]"]),
    ("major-decisions", "Major Decisions", "Section 06",
     ["[Decision 1 — date — decision maker — rationale]",
      "[Decision 2 — date — decision maker — rationale]",
      "[Decision 3 — date — decision maker — rationale]",
      "[Open decisions awaiting resolution]"]),
    ("risk-overview", "Risk Overview", "Section 07",
     ["[Risk 1 — likelihood / impact — mitigation — owner]",
      "[Risk 2 — likelihood / impact — mitigation — owner]",
      "[Risk 3 — likelihood / impact — mitigation — owner]",
      "[Escalation path for new risks]"]),
    ("stakeholder-review", "Stakeholder Review", "Section 08",
     ["[Stakeholder / team — role in the plan — engagement level]",
      "[Key contacts and communication cadence]",
      "[Outstanding stakeholder concerns]",
      "[Approvals and sign-offs received]"]),
    ("immediate-priorities", "Immediate Priorities", "Section 09",
     ["[Priority 1 — owner — due date]",
      "[Priority 2 — owner — due date]",
      "[Priority 3 — owner — due date]",
      "[First 30 / 60 / 90 day focus]"]),
    ("lessons-learned", "Lessons Learned", "Section 10",
     ["[What worked well and should be repeated]",
      "[What did not work and why]",
      "[Process improvements recommended]",
      "[Advice for the next plan transition]"]),
]

# ----------------------------------------------------------------- styles ---

STYLES = {
    "1-boardroom-blue": dict(
        name="Boardroom Blue", mode="clean",
        bg=(255, 255, 255), panel=(11, 37, 69), band=(11, 37, 69),
        ink=(11, 37, 69), sub=(90, 105, 125), accent=(31, 111, 191),
        accent2=(140, 178, 217), light=(226, 236, 246),
        head_font="Calibri", body_font="Calibri", title_on_dark=True,
    ),
    "2-hand-drawn": dict(
        name="Hand-Drawn Sketch", mode="sketch",
        bg=(250, 247, 240), panel=(250, 247, 240), band=(37, 44, 59),
        ink=(37, 44, 59), sub=(110, 112, 108), accent=(196, 90, 42),
        accent2=(122, 140, 105), light=(235, 229, 214),
        head_font="Georgia", body_font="Calibri", title_on_dark=False,
    ),
    "3-executive-dark": dict(
        name="Executive Dark", mode="clean",
        bg=(18, 20, 24), panel=(24, 27, 33), band=(24, 27, 33),
        ink=(238, 234, 224), sub=(150, 152, 158), accent=(198, 164, 92),
        accent2=(120, 104, 70), light=(40, 44, 52),
        head_font="Calibri Light", body_font="Calibri", title_on_dark=True,
    ),
    "4-horizon-teal": dict(
        name="Horizon Teal", mode="clean",
        bg=(246, 250, 250), panel=(6, 74, 82), band=(6, 74, 82),
        ink=(10, 42, 46), sub=(84, 110, 112), accent=(0, 128, 128),
        accent2=(94, 191, 178), light=(210, 235, 232),
        head_font="Calibri", body_font="Calibri", title_on_dark=True,
    ),
    "5-minimal-editorial": dict(
        name="Minimal Editorial", mode="clean",
        bg=(247, 244, 239), panel=(247, 244, 239), band=(45, 41, 38),
        ink=(45, 41, 38), sub=(120, 112, 104), accent=(170, 74, 51),
        accent2=(150, 140, 125), light=(230, 224, 214),
        head_font="Georgia", body_font="Calibri", title_on_dark=False,
    ),
}

# ------------------------------------------------------------ PIL helpers ---

random.seed(42)


def _jit(p, amt):
    return (p[0] + random.uniform(-amt, amt), p[1] + random.uniform(-amt, amt))


def sk_line(d, p1, p2, color, w, mode, amt=3.0):
    """Line; in sketch mode draw as jittered polyline with a double stroke."""
    if mode != "sketch":
        d.line([p1, p2], fill=color, width=w)
        return
    n = max(3, int(math.dist(p1, p2) / 28))
    for _pass in range(2):
        pts = []
        for i in range(n + 1):
            t = i / n
            x = p1[0] + (p2[0] - p1[0]) * t
            y = p1[1] + (p2[1] - p1[1]) * t
            pts.append(_jit((x, y), amt if 0 < i < n else amt * 0.4))
        d.line(pts, fill=color, width=max(1, w - _pass), joint="curve")


def sk_circle(d, c, r, color, w, mode, fill=None, amt=3.0):
    box = [c[0] - r, c[1] - r, c[0] + r, c[1] + r]
    if mode != "sketch":
        d.ellipse(box, outline=color, width=w, fill=fill)
        return
    if fill:
        d.ellipse(box, fill=fill)
    for _pass in range(2):
        pts = []
        start = random.uniform(0, 6.28)
        for i in range(37):
            a = start + i / 36 * 2 * math.pi
            rr = r + random.uniform(-amt, amt)
            pts.append((c[0] + rr * math.cos(a), c[1] + rr * math.sin(a)))
        d.line(pts, fill=color, width=max(1, w - _pass), joint="curve")


def sk_rect(d, box, color, w, mode, fill=None, amt=3.0):
    x0, y0, x1, y1 = box
    if mode != "sketch":
        d.rectangle(box, outline=color, width=w, fill=fill)
        return
    if fill:
        d.rectangle(box, fill=fill)
    for a, b in [((x0, y0), (x1, y0)), ((x1, y0), (x1, y1)),
                 ((x1, y1), (x0, y1)), ((x0, y1), (x0, y0))]:
        sk_line(d, a, b, color, w, mode, amt)


def sk_arrow(d, p1, p2, color, w, mode):
    sk_line(d, p1, p2, color, w, mode)
    ang = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
    L = 26
    for da in (2.6, -2.6):
        q = (p2[0] + L * math.cos(ang + da), p2[1] + L * math.sin(ang + da))
        sk_line(d, p2, q, color, w, mode, amt=1.5)


def hatch(d, box, color, mode, step=18, w=3):
    """Diagonal hatch fill (sketchy shading)."""
    x0, y0, x1, y1 = box
    span = int((x1 - x0) + (y1 - y0))
    for s in range(0, span, step):
        a = (x0 + s, y0) if s < (x1 - x0) else (x1, y0 + (s - (x1 - x0)))
        b = (x0, y0 + s) if s < (y1 - y0) else (x0 + (s - (y1 - y0)), y1)
        sk_line(d, a, b, color, w, mode, amt=1.5)


# ------------------------------------------------------------------ motifs --

S = 900  # canvas size for square motifs


def _canvas(w=S, h=S):
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    return img, ImageDraw.Draw(img)


def motif_hero(st):
    """Title-slide composition: rising arcs + horizon line + sun disc."""
    img, d = _canvas(1200, 1200)
    mode, acc, acc2, ink = st["mode"], st["accent"], st["accent2"], st["ink"]
    cx, cy = 600, 760
    for i, r in enumerate(range(160, 561, 80)):
        col = acc if i % 2 == 0 else acc2
        if mode == "sketch":
            sk_circle(d, (cx, cy), r, col, 6, mode, amt=4)
        else:
            d.arc([cx - r, cy - r, cx + r, cy + r], 180, 360, fill=col, width=10)
    sk_line(d, (60, cy), (1140, cy), ink, 8, mode)
    sk_circle(d, (cx, cy - 300), 56, acc, 7, mode,
              fill=acc + (60,) if mode != "sketch" else None)
    for x in range(140, 1101, 120):
        sk_line(d, (x, cy + 60), (x, cy + 60 + (18 if (x // 120) % 2 else 34)),
                st["sub"], 5, mode)
    return img


def motif_rings(st):        # executive overview
    img, d = _canvas()
    mode = st["mode"]
    for r, col, w in [(330, st["accent2"], 6), (250, st["accent"], 8),
                      (170, st["ink"], 6), (90, st["accent"], 8)]:
        sk_circle(d, (450, 430), r, col, w, mode, amt=4)
    sk_circle(d, (450, 430), 26, st["accent"], 6, mode, fill=st["accent"])
    sk_line(d, (450, 430), (760, 160), st["sub"], 5, mode)
    sk_circle(d, (760, 160), 14, st["ink"], 5, mode, fill=st["ink"])
    return img


def motif_gauge(st):        # project status
    img, d = _canvas()
    mode, c = st["mode"], (450, 560)
    if mode == "sketch":
        sk_circle(d, c, 300, st["ink"], 7, mode, amt=4)
    else:
        d.arc([150, 260, 750, 860], 180, 360, fill=st["ink"], width=12)
        d.arc([210, 320, 690, 800], 180, 305, fill=st["accent"], width=22)
    for i in range(7):
        a = math.pi + i * math.pi / 6
        p1 = (c[0] + 268 * math.cos(a), c[1] + 268 * math.sin(a))
        p2 = (c[0] + 300 * math.cos(a), c[1] + 300 * math.sin(a))
        sk_line(d, p1, p2, st["sub"], 6, mode)
    a = math.pi * 1.72
    sk_arrow(d, c, (c[0] + 230 * math.cos(a), c[1] + 230 * math.sin(a)),
             st["accent"], 9, mode)
    sk_circle(d, c, 22, st["ink"], 6, mode, fill=st["ink"])
    return img


def motif_frames(st):       # project scope
    img, d = _canvas()
    mode = st["mode"]
    sk_rect(d, (140, 140, 760, 760), st["accent2"], 6, mode, amt=4)
    sk_rect(d, (220, 220, 680, 680), st["accent"], 8, mode, amt=4)
    sk_rect(d, (300, 300, 600, 600), st["ink"], 6, mode, amt=3)
    if mode == "sketch":
        hatch(d, (300, 300, 600, 600), st["accent"], mode, step=26, w=3)
    else:
        d.rectangle((300, 300, 600, 600), fill=st["accent"] + (40,))
    sk_line(d, (600, 300), (780, 170), st["sub"], 5, mode)
    sk_circle(d, (790, 162), 12, st["accent"], 5, mode, fill=st["accent"])
    return img


def motif_timeline(st):     # milestones
    img, d = _canvas()
    mode = st["mode"]
    sk_line(d, (90, 640), (810, 260), st["ink"], 8, mode)
    for i, t in enumerate([0.08, 0.36, 0.64, 0.92]):
        x = 90 + t * 720
        y = 640 - t * 380
        filled = i < 2
        sk_circle(d, (x, y), 34, st["accent"] if filled else st["accent2"], 8,
                  mode, fill=(st["accent"] if filled and mode != "sketch" else None))
        if filled and mode == "sketch":
            hatch(d, (x - 24, y - 24, x + 24, y + 24), st["accent"], mode, 10, 3)
        sk_line(d, (x, y + 46), (x, y + 96), st["sub"], 5, mode)
    sk_arrow(d, (810, 260), (860, 234), st["ink"], 8, mode)
    return img


def motif_lanes(st):        # work stream status
    img, d = _canvas()
    mode = st["mode"]
    fills = [0.85, 0.55, 0.7, 0.35]
    for i, f in enumerate(fills):
        y0 = 190 + i * 150
        sk_rect(d, (110, y0, 790, y0 + 84), st["sub"], 6, mode, amt=3)
        x1 = 110 + f * 680
        if mode == "sketch":
            hatch(d, (116, y0 + 6, x1, y0 + 78),
                  st["accent"] if i % 2 == 0 else st["accent2"], mode, 20, 3)
        else:
            d.rectangle((116, y0 + 6, x1, y0 + 78),
                        fill=(st["accent"] if i % 2 == 0 else st["accent2"]) + (200,))
        sk_circle(d, (830, y0 + 42), 16, st["accent"] if f > 0.6 else st["accent2"],
                  6, mode, fill=None)
    return img


def motif_fork(st):         # major decisions
    img, d = _canvas()
    mode = st["mode"]
    sk_line(d, (150, 450), (430, 450), st["ink"], 9, mode)
    sk_circle(d, (450, 450), 40, st["accent"], 8, mode)
    for dy, col in [(-230, st["accent"]), (0, st["ink"]), (230, st["accent2"])]:
        sk_arrow(d, (486, 450 + dy * 0.12), (760, 450 + dy), col, 8, mode)
        sk_circle(d, (790, 450 + dy), 22, col, 7, mode,
                  fill=col if mode != "sketch" else None)
    if mode == "sketch":
        hatch(d, (426, 426, 474, 474), st["accent"], mode, 12, 3)
    return img


def motif_matrix(st):       # risk overview
    img, d = _canvas()
    mode = st["mode"]
    sk_rect(d, (150, 150, 750, 750), st["ink"], 7, mode, amt=4)
    sk_line(d, (450, 150), (450, 750), st["sub"], 5, mode)
    sk_line(d, (150, 450), (750, 450), st["sub"], 5, mode)
    if mode == "sketch":
        hatch(d, (455, 155, 745, 445), st["accent"], mode, 24, 3)
    else:
        d.rectangle((455, 155, 745, 445), fill=st["accent"] + (45,))
    for (x, y, r, col) in [(280, 620, 26, st["accent2"]), (360, 320, 30, st["sub"]),
                           (600, 560, 26, st["accent2"]), (610, 260, 38, st["accent"]),
                           (540, 350, 24, st["accent"])]:
        sk_circle(d, (x, y), r, col, 7, mode,
                  fill=col if mode != "sketch" else None)
    sk_arrow(d, (150, 780), (750, 780), st["ink"], 6, mode)
    sk_arrow(d, (120, 750), (120, 150), st["ink"], 6, mode)
    return img


def motif_network(st):      # stakeholder review
    img, d = _canvas()
    mode = st["mode"]
    hub = (450, 450)
    nodes = [(450, 150, 34), (740, 330, 30), (700, 680, 26),
             (240, 700, 30), (160, 350, 26), (620, 500, 22)]
    for (x, y, r) in nodes:
        sk_line(d, hub, (x, y), st["sub"], 5, mode)
    sk_line(d, nodes[0][:2], nodes[1][:2], st["accent2"], 4, mode)
    sk_line(d, nodes[3][:2], nodes[4][:2], st["accent2"], 4, mode)
    sk_circle(d, hub, 52, st["accent"], 9, mode,
              fill=st["accent"] if mode != "sketch" else None)
    if mode == "sketch":
        hatch(d, (415, 415, 485, 485), st["accent"], mode, 12, 3)
    for i, (x, y, r) in enumerate(nodes):
        col = st["ink"] if i % 2 == 0 else st["accent2"]
        sk_circle(d, (x, y), r, col, 7, mode)
    return img


def motif_stack(st):        # immediate priorities
    img, d = _canvas()
    mode = st["mode"]
    for i in range(3):
        y = 620 - i * 190
        wdt = 520 - i * 120
        x0 = 450 - wdt / 2
        col = [st["accent2"], st["sub"], st["accent"]][i]
        sk_rect(d, (x0, y, x0 + wdt, y + 120), col, 8, mode, amt=3)
        if i == 2:
            if mode == "sketch":
                hatch(d, (x0 + 6, y + 6, x0 + wdt - 6, y + 114), col, mode, 22, 3)
            else:
                d.rectangle((x0 + 6, y + 6, x0 + wdt - 6, y + 114),
                            fill=col + (50,))
    sk_arrow(d, (450, 200), (450, 110), st["accent"], 9, mode)
    return img


def motif_bulb(st):         # lessons learned
    img, d = _canvas()
    mode, c = st["mode"], (450, 400)
    sk_circle(d, c, 180, st["accent"], 9, mode)
    sk_rect(d, (400, 580, 500, 650), st["ink"], 7, mode, amt=2)
    sk_line(d, (410, 690), (490, 690), st["ink"], 7, mode)
    for i in range(8):
        a = -math.pi / 2 + (i - 3.5) * 0.42
        p1 = (c[0] + 215 * math.cos(a), c[1] + 215 * math.sin(a))
        p2 = (c[0] + 290 * math.cos(a), c[1] + 290 * math.sin(a))
        sk_line(d, p1, p2, st["accent2"], 6, mode)
    sk_line(d, (405, 470), (440, 530), st["sub"], 6, mode)
    sk_line(d, (495, 470), (460, 530), st["sub"], 6, mode)
    sk_line(d, (440, 530), (460, 530), st["sub"], 6, mode)
    return img


MOTIFS = {
    "hero": motif_hero,
    "executive-overview": motif_rings,
    "project-status": motif_gauge,
    "project-scope": motif_frames,
    "milestones": motif_timeline,
    "work-stream-status": motif_lanes,
    "major-decisions": motif_fork,
    "risk-overview": motif_matrix,
    "stakeholder-review": motif_network,
    "immediate-priorities": motif_stack,
    "lessons-learned": motif_bulb,
}


def art_path(style_key, motif_key):
    st = STYLES[style_key]
    path = os.path.join(ART, f"{style_key}--{motif_key}.png")
    random.seed(hash((style_key, motif_key)) % (2**31))
    img = MOTIFS[motif_key](st)
    if st["mode"] == "sketch":
        img = img.filter(ImageFilter.GaussianBlur(0.6))
    img.save(path)
    return path


def gradient_bg(style_key, top, bottom):
    path = os.path.join(ART, f"{style_key}--bg.png")
    img = Image.new("RGB", (1600, 900))
    for y in range(900):
        t = y / 899
        img.paste(tuple(int(top[i] + (bottom[i] - top[i]) * t) for i in range(3)),
                  (0, y, 1600, y + 1))
    img.save(path)
    return path


# ------------------------------------------------------------ pptx helpers --


def add_text(slide, x, y, w, h, text, size, color, font="Calibri", bold=False,
             align=PP_ALIGN.LEFT, italic=False, spacing=None, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = RGBColor(*color)
    if spacing is not None:
        p.line_spacing = spacing
    return tb


def add_bullets(slide, x, y, w, h, items, size, color, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = "•  " + item
        f = r.font
        f.size = Pt(size)
        f.name = font
        f.color.rgb = RGBColor(*color)
    return tb


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=None,
            shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor(*fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor(*line)
        sp.line.width = Pt(line_w or 1)
    return sp


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*color)


# --------------------------------------------------------------- builders ---


def build_deck(style_key):
    st = STYLES[style_key]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    grad = None
    if style_key == "4-horizon-teal":
        grad = gradient_bg(style_key, (246, 250, 250), (208, 232, 229))

    # ---- slide 1: title -----------------------------------------------
    s = prs.slides.add_slide(blank)
    set_bg(s, st["bg"])
    if grad:
        s.shapes.add_picture(grad, 0, 0, SLIDE_W, SLIDE_H)

    dark_title = st["title_on_dark"]
    if dark_title:
        add_box(s, 0, 0, SLIDE_W, SLIDE_H, fill=st["band"])
    # hero art on the right
    s.shapes.add_picture(art_path(style_key, "hero"),
                         Inches(7.6), Inches(0.9), height=Inches(5.7))
    tcol = (255, 255, 255) if dark_title else st["ink"]
    scol = st["accent2"] if dark_title else st["accent"]
    add_box(s, Inches(0.9), Inches(2.05), Inches(1.6), Pt(4), fill=st["accent"])
    add_text(s, Inches(0.85), Inches(2.3), Inches(6.6), Inches(1.9),
             "Executive Handoff Summary", 44, tcol, st["head_font"], bold=True,
             spacing=1.02)
    add_text(s, Inches(0.9), Inches(4.15), Inches(6.2), Inches(0.6),
             "[Plan Name]", 24, scol, st["head_font"], italic=True)
    add_text(s, Inches(0.9), Inches(4.85), Inches(6.2), Inches(0.6),
             "Overview of the Corporate Retirement Plan Process", 16,
             (225, 225, 225) if dark_title else st["sub"], st["body_font"])
    add_text(s, Inches(0.9), Inches(6.6), Inches(7.0), Inches(0.4),
             "[Presenter Name]  •  [Date]  •  Transamerica", 12,
             (200, 200, 200) if dark_title else st["sub"], st["body_font"])

    # ---- content slides -------------------------------------------------
    for i, (key, title, kicker, bullets) in enumerate(SECTIONS, start=2):
        s = prs.slides.add_slide(blank)
        set_bg(s, st["bg"])
        if grad:
            s.shapes.add_picture(grad, 0, 0, SLIDE_W, SLIDE_H)

        # style-specific chrome
        if style_key == "1-boardroom-blue":
            add_box(s, 0, 0, Inches(0.28), SLIDE_H, fill=st["band"])
            add_box(s, 0, 0, SLIDE_W, Inches(0.1), fill=st["accent"])
        elif style_key == "2-hand-drawn":
            add_box(s, Inches(0.45), Inches(0.42), Inches(12.45), Inches(6.33),
                    line=st["ink"], line_w=1.5)
        elif style_key == "3-executive-dark":
            add_box(s, 0, Inches(1.62), Inches(13.333), Pt(1.4), fill=st["accent2"])
            add_box(s, Inches(0.9), Inches(1.6), Inches(2.2), Pt(3), fill=st["accent"])
        elif style_key == "4-horizon-teal":
            add_box(s, 0, Inches(7.28), SLIDE_W, Inches(0.22), fill=st["accent"])
            add_box(s, Inches(0.9), Inches(1.62), Inches(1.4), Pt(4), fill=st["accent"])
        elif style_key == "5-minimal-editorial":
            add_box(s, Inches(0.9), Inches(1.66), Inches(11.53), Pt(1.2),
                    fill=st["ink"])
            add_box(s, Inches(0.9), Inches(1.62), Inches(1.0), Pt(4),
                    fill=st["accent"])

        kx = Inches(0.9)
        add_text(s, kx, Inches(0.55), Inches(4), Inches(0.4),
                 kicker.upper(), 12, st["accent"], st["body_font"], bold=True)
        add_text(s, kx, Inches(0.85), Inches(8.6), Inches(0.9),
                 title, 32, st["ink"], st["head_font"], bold=True)

        # motif art, right side
        s.shapes.add_picture(art_path(style_key, key),
                             Inches(8.85), Inches(2.15), height=Inches(4.4))

        # body placeholder panel + bullets
        if style_key == "2-hand-drawn":
            add_box(s, Inches(0.9), Inches(2.1), Inches(7.5), Inches(4.5),
                    line=st["sub"], line_w=1.0)
        elif style_key == "3-executive-dark":
            add_box(s, Inches(0.9), Inches(2.1), Inches(7.5), Inches(4.5),
                    fill=st["light"])
        elif style_key == "1-boardroom-blue":
            add_box(s, Inches(0.9), Inches(2.1), Inches(7.5), Inches(4.5),
                    fill=st["light"])
        elif style_key == "4-horizon-teal":
            add_box(s, Inches(0.9), Inches(2.1), Inches(7.5), Inches(4.5),
                    fill=(255, 255, 255), line=st["accent2"], line_w=1.0)
        else:
            add_box(s, Inches(0.9), Inches(2.1), Inches(7.5), Inches(4.5),
                    fill=(255, 255, 255), line=st["light"], line_w=1.0)

        add_bullets(s, Inches(1.25), Inches(2.45), Inches(6.9), Inches(3.9),
                    bullets, 15, st["ink"], st["body_font"])

        # footer
        add_text(s, Inches(0.9), Inches(6.95), Inches(6), Inches(0.35),
                 "[Plan Name]  —  Executive Handoff Summary", 10,
                 st["sub"], st["body_font"])
        add_text(s, Inches(12.1), Inches(6.95), Inches(0.9), Inches(0.35),
                 str(i), 10, st["sub"], st["body_font"], align=PP_ALIGN.RIGHT)

    out = os.path.join(OUT, f"executive-handoff-template-{style_key}.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    for k in STYLES:
        print("built:", build_deck(k))
